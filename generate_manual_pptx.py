# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

BLUE_DARK  = RGBColor(0x0D, 0x3C, 0x6E)
BLUE_MID   = RGBColor(0x1A, 0x6F, 0xB5)
BLUE_LIGHT = RGBColor(0xD6, 0xE8, 0xF7)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_LIGHT = RGBColor(0xF0, 0xF4, 0xF8)
ORANGE     = RGBColor(0xE8, 0x7A, 0x1E)
GREEN      = RGBColor(0x1E, 0x8A, 0x44)
TEXT_DARK  = RGBColor(0x1A, 0x1A, 0x2E)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ── helpers ───────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=None, line=None, lw=Pt(0)):
    from pptx.enum.shapes import PP_PLACEHOLDER
    shp = slide.shapes.add_shape(1, l, t, w, h)
    shp.line.width = lw
    if fill:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line:
        shp.line.color.rgb = line
    else:
        shp.line.fill.background()
    return shp

def tb(slide, text, l, t, w, h, sz=Pt(14), bold=False, clr=None, align=PP_ALIGN.LEFT):
    if clr is None: clr = TEXT_DARK
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame; tf.word_wrap = True
    first = True
    for line in text.split('\n'):
        if first: p = tf.paragraphs[0]; first = False
        else: p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = sz; run.font.bold = bold
        run.font.color.rgb = clr; run.font.name = u'\u5fae\u8f6f\u96c5\u9ed1'
    return txb

def blue_bg(slide):
    add_rect(slide, 0, 0, W, H, fill=BLUE_DARK)
    gc = RGBColor(0x1E, 0x5A, 0x9A)
    for i in range(1, 14):
        y = int(H * i / 14)
        r = slide.shapes.add_shape(1, 0, y, W, Emu(2500))
        r.fill.solid(); r.fill.fore_color.rgb = gc; r.line.fill.background()
    for i in range(1, 24):
        x = int(W * i / 24)
        r = slide.shapes.add_shape(1, x, 0, Emu(2500), H)
        r.fill.solid(); r.fill.fore_color.rgb = gc; r.line.fill.background()

def cbg(slide):
    add_rect(slide, 0, 0, W, H, fill=GRAY_LIGHT)
    add_rect(slide, 0, 0, W, Inches(0.75), fill=BLUE_DARK)
    add_rect(slide, 0, H-Inches(0.3), W, Inches(0.3), fill=BLUE_MID)
    add_rect(slide, 0, Inches(0.75), Inches(0.12), H-Inches(1.05), fill=BLUE_MID)

def ph(slide, l, t, w, h, cap):
    add_rect(slide, l, t, w, h, fill=BLUE_LIGHT, line=BLUE_MID, lw=Pt(1.5))
    tb(slide, cap, l, t+h//2-Inches(0.22), w, Inches(0.44),
       sz=Pt(11), clr=BLUE_MID, align=PP_ALIGN.CENTER)

def callout(slide, text, l, t, w, h, bg=ORANGE, fg=WHITE, sz=Pt(13)):
    add_rect(slide, l, t, w, h, fill=bg)
    tb(slide, text, l+Inches(0.12), t+Inches(0.07), w-Inches(0.24), h-Inches(0.14),
       sz=sz, clr=fg)

def steps(slide, items, top=Inches(1.7), left=Inches(0.7), sz=Pt(14)):
    txb = slide.shapes.add_textbox(left, top, Inches(12.0), Inches(5.5))
    tf  = txb.text_frame; tf.word_wrap = True
    first = True
    for item in items:
        if first: p = tf.paragraphs[0]; first = False
        else: p = tf.add_paragraph()
        p.space_before = Pt(5)
        run = p.add_run(); run.text = item
        run.font.size = sz; run.font.color.rgb = TEXT_DARK
        run.font.name = u'\u5fae\u8f6f\u96c5\u9ed1'

def hdr(slide, text):
    tb(slide, text, Inches(0.22), Inches(0.1), Inches(12), Inches(0.55),
       sz=Pt(20), bold=True, clr=WHITE)

# ── SLIDE 1 Cover ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
blue_bg(s)
tb(s, 'CADCheckTool', Inches(1), Inches(1.6), Inches(11), Inches(1.1),
   sz=Pt(52), bold=True, clr=WHITE, align=PP_ALIGN.CENTER)
tb(s, u'\u7528\u6237\u7aef\u64cd\u4f5c\u624b\u518c',  # 用户端操作手册
   Inches(1), Inches(2.7), Inches(11), Inches(0.85),
   sz=Pt(36), bold=True, clr=RGBColor(0xA8,0xD4,0xFF), align=PP_ALIGN.CENTER)
tb(s, u'AutoCAD \u5de5\u7a0b\u56fe\u7eb8\u81ea\u52a8\u5316\u68c0\u67e5\u63d2\u4ef6  \uff5c  \u9762\u5411\u6700\u7ec8\u7528\u6237',
   Inches(1), Inches(3.55), Inches(11), Inches(0.5),
   sz=Pt(16), clr=RGBColor(0xB0,0xC8,0xE8), align=PP_ALIGN.CENTER)
ph(s, Inches(5.4), Inches(4.3), Inches(2.5), Inches(1.6),
   u'\u3010\u5360\u4f4d\uff1a\u4ea7\u54c1 Logo / \u5c01\u9762\u56fe\u3011')
tb(s, u'\u7248\u672c V1.0  \u00b7  2024',
   Inches(1), Inches(6.7), Inches(11), Inches(0.35),
   sz=Pt(12), clr=RGBColor(0x80,0xA8,0xD0), align=PP_ALIGN.CENTER)

# ── SLIDE 2 TOC ───────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
cbg(s)
tb(s, u'\u76ee  \u5f55', Inches(0.22), Inches(0.1), Inches(8), Inches(0.55),
   sz=Pt(20), bold=True, clr=WHITE)
toc = [
    (u'01', u'\u4f7f\u7528\u987b\u77e5 & \u6ce8\u610f\u4e8b\u9879'),
    (u'02', u'\u9996\u6b21\u5b89\u88c5\u4e0e\u914d\u7f6e'),
    (u'03', u'\u5feb\u6377\u547d\u4ee4\u5165\u53e3'),
    (u'04', u'\u5355\u5f20\u68c0\u67e5\u2014\u2014\u64cd\u4f5c\u6d41\u7a0b'),
    (u'05', u'\u5355\u5f20\u68c0\u67e5\u2014\u2014\u7f16\u8f91\u4e0e\u590d\u67e5'),
    (u'06', u'\u6279\u91cf\u68c0\u67e5\u2014\u2014\u6267\u884c\u6d41\u7a0b'),
    (u'07', u'\u6279\u91cf\u68c0\u67e5\u2014\u2014\u67e5\u770b\u62a5\u544a'),
    (u'08', u'\u6279\u91cf\u6ce8\u91ca\u6e05\u9664'),
    (u'09', u'\u53ef\u68c0\u67e5\u95ee\u9898\u6c47\u603b'),
    (u'10', u'\u62a5\u9519\u4e0e\u9632\u5d29\u6e83\u673a\u5236'),
    (u'11', u'\u8f6f\u4ef6\u5378\u8f7d'),
    (u'12', u'\u5feb\u901f\u6838\u67e5\u6e05\u5355 & \u652f\u6301'),
]
cw = Inches(5.8)
for i, (num, title) in enumerate(toc):
    col = i % 2; row = i // 2
    lx = Inches(0.6) + col*cw
    ty = Inches(1.05) + row*Inches(0.72)
    add_rect(s, lx, ty+Inches(0.05), Inches(0.5), Inches(0.44), fill=BLUE_MID)
    tb(s, num, lx, ty+Inches(0.04), Inches(0.5), Inches(0.46),
       sz=Pt(11), bold=True, clr=WHITE, align=PP_ALIGN.CENTER)
    tb(s, title, lx+Inches(0.58), ty, Inches(5.1), Inches(0.55), sz=Pt(15), clr=TEXT_DARK)

# ── SLIDE 3 Before you start ──────────────────────────────────
s = prs.slides.add_slide(BLANK)
cbg(s)
hdr(s, u'\u4f7f\u7528\u987b\u77e5  &  \u6ce8\u610f\u4e8b\u9879')
callout(s, u'\u26a0  \u5b89\u88c5\u524d\u5fc5\u987b\u5173\u95ed AutoCAD',
        Inches(0.6), Inches(0.9), Inches(11.8), Inches(0.65), sz=Pt(15))
callout(s, u'\u26a0  \u6267\u884c\u6279\u91cf\u68c0\u67e5 / \u6279\u91cf\u6ce8\u91ca\u6e05\u9664\u65f6\uff0c\u76ee\u6807\u56fe\u7eb8\u6587\u4ef6\u5fc5\u987b\u5904\u4e8e\u5173\u95ed\u72b6\u6001',
        Inches(0.6), Inches(1.65), Inches(11.8), Inches(0.65), sz=Pt(15))
tb(s, u'\u5176\u4ed6\u4f7f\u7528\u5efa\u8bae', Inches(0.7), Inches(2.55), Inches(5), Inches(0.4),
   sz=Pt(16), bold=True, clr=BLUE_DARK)
tips = [
    u'\u2714  \u9996\u6b21\u5b89\u88c5\u540e\uff0c\u540e\u7eed\u6253\u5f00 AutoCAD \u5c06\u81ea\u52a8\u52a0\u8f7d\u63d2\u4ef6\uff0c\u65e0\u9700\u91cd\u590d\u4e0b\u8f7d\u5b89\u88c5\u3002',
    u'\u2714  \u5f39\u7a97\u68c0\u67e5\u7ed3\u679c\u4ec5\u4f9b\u53c2\u8003\uff0c\u56fe\u7eb8\u5185\u7684\u6807\u6ce8\u5185\u5bb9\u4e3a\u6700\u7ec8\u6743\u5a01\u4f9d\u636e\u3002',
    u'\u2714  \u82e5\u63d2\u4ef6\u51fa\u73b0\u4e25\u91cd\u62a5\u9519\uff0cAutoCAD \u53ef\u80fd\u5173\u95ed\uff0c\u4f46\u539f\u59cb DWG \u6587\u4ef6\u4e0d\u4f1a\u635f\u574f\u3002',
    u'\u2714  \u5efa\u8bae\u5728\u6267\u884c\u6279\u91cf\u68c0\u67e5\u524d\uff0c\u786e\u8ba4\u6240\u6709\u76ee\u6807\u56fe\u7eb8\u5df2\u4fdd\u5b58\u5e76\u5173\u95ed\u3002',
]
steps(s, tips, top=Inches(3.05), left=Inches(0.7), sz=Pt(14))

# ── SLIDE 4 Install ───────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
cbg(s)
hdr(s, u'\u9996\u6b21\u5b89\u88c5\u4e0e\u914d\u7f6e')
si = [
    u'\u2460 \u786e\u8ba4 AutoCAD \u5df2\u5b8c\u5168\u5173\u95ed\uff08\u5305\u62ec\u540e\u53f0\u8fdb\u7a0b\uff09',
    u'\u2461 \u4e0b\u8f7d\u5b89\u88c5\u5305\uff0c\u53cc\u51fb\u8fd0\u884c setup.exe',
    u'   \u279e \u63a8\u8350\uff1a\u53f3\u952e \u2192 \u4ee5\u7ba1\u7406\u5458\u8eab\u4efd\u8fd0\u884c',
    u'\u2462 \u5728\u5b89\u88c5\u754c\u9762\u4e2d\u70b9\u51fb [Install] \u6309\u9215',
    u'\u2463 \u7b49\u5f85\u5b89\u88c5\u5b8c\u6210',
    u'   \u2714  \u5b8c\u6210\u540e\u5373\u914d\u7f6e\u6210\u529f',
    u'   \u2714  \u540e\u7eed\u542f\u52a8 AutoCAD \u5c06\u81ea\u52a8\u52a0\u8f7d\u63d2\u4ef6\uff0c\u65e0\u9700\u91cd\u590d\u5b89\u88c5',
]
steps(s, si, top=Inches(0.95), sz=Pt(15))
ph(s, Inches(7.6), Inches(1.0), Inches(4.8), Inches(3.2),
   u'\u3010\u622a\u56fe\u5360\u4f4d\uff1a\u5b89\u88c5\u754c\u9762 (setup.exe)\u3011')
callout(s, u'[tip] \u5c0f\u63d0\u793a\uff1a\u5b89\u88c5\u5b8c\u6210\u540e\uff0c\u4e0b\u6b21\u6253\u5f00 AutoCAD \u5373\u53ef\u76f4\u63a5\u4f7f\u7528\uff0c\u65e0\u9700\u4efb\u4f55\u989d\u5916\u64cd\u4f5c\u3002',
        Inches(0.6), Inches(6.35), Inches(11.8), Inches(0.65),
        bg=GREEN, fg=WHITE, sz=Pt(13))

# ── SLIDE 5 Command ───────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
cbg(s)
hdr(s, u'\u5feb\u6377\u547d\u4ee4\u5165\u53e3')
tb(s, u'\u6253\u5f00\u63d2\u4ef6\u7a97\u53e3\u7684\u65b9\u5f0f', Inches(0.7), Inches(0.95), Inches(8), Inches(0.5),
   sz=Pt(17), bold=True, clr=BLUE_DARK)
add_rect(s, Inches(3.5), Inches(1.65), Inches(6.3), Inches(1.1), fill=BLUE_DARK)
tb(s, 'CHECKDRAWING', Inches(3.5), Inches(1.68), Inches(6.3), Inches(1.1),
   sz=Pt(34), bold=True, clr=RGBColor(0xA8,0xD4,0xFF), align=PP_ALIGN.CENTER)
tb(s, u'\u5728 AutoCAD \u547d\u4ee4\u884c\u4e2d\u8f93\u5165\u4ee5\u4e0a\u547d\u4ee4\uff0c\u7136\u540e\u6309 Enter \u952e\uff0c\n\u5373\u53ef\u6253\u5f00 CADCheckTool \u63d2\u4ef6\u7a97\u53e3\u3002',
   Inches(0.7), Inches(1.7), Inches(2.6), Inches(1.4), sz=Pt(14), clr=TEXT_DARK)
ph(s, Inches(0.7), Inches(3.2), Inches(11.8), Inches(3.4),
   u'\u3010\u622a\u56fe\u5360\u4f4d\uff1aAutoCAD \u547d\u4ee4\u884c\u8f93\u5165 CHECKDRAWING / \u63d2\u4ef6\u4e3b\u7a97\u53e3\u3011')

# ── SLIDE 6 Single check overview ─────────────────────────────
s = prs.slides.add_slide(BLANK)
cbg(s)
hdr(s, u'\u5355\u5f20\u68c0\u67e5\u2014\u2014\u64cd\u4f5c\u6982\u89c8')
flow = [
    (u'\u70b9\u51fb\n\u300c\u5355\u5f20\u68c0\u67e5\u300d', Inches(0.6)),
    (u'\u70b9\u51fb\n\u300c\u68c0\u67e5\u5f53\u524d\u56fe\u7eb8\u300d', Inches(3.2)),
    (u'\u7b49\u5f85\u68c0\u67e5\u5b8c\u6210', Inches(5.9)),
    (u'\u67e5\u770b\u5f39\u7a97\n+\u56fe\u7eb8\u6807\u6ce8', Inches(8.5)),
]
for label, lx in flow:
    add_rect(s, lx, Inches(2.1), Inches(2.3), Inches(1.0), fill=BLUE_MID)
    tb(s, label, lx, Inches(2.1), Inches(2.3), Inches(1.0),
       sz=Pt(13), bold=True, clr=WHITE, align=PP_ALIGN.CENTER)
for ax in [Inches(2.9), Inches(5.5), Inches(8.2)]:
    tb(s, u'\u25b6', ax, Inches(2.35), Inches(0.35), Inches(0.5),
       sz=Pt(18), clr=BLUE_MID, align=PP_ALIGN.CENTER)
ph(s, Inches(0.6), Inches(3.5), Inches(5.6), Inches(3.0),
   u'\u3010\u622a\u56fe\u5360\u4f4d\uff1a\u5355\u5f20\u68c0\u67e5\u5165\u53e3\u6309\u9215\u3011')
ph(s, Inches(6.5), Inches(3.5), Inches(6.2), Inches(3.0),
   u'\u3010\u622a\u56fe\u5360\u4f4d\uff1a\u68c0\u67e5\u5b8c\u6210\u5f39\u7a97 & \u56fe\u7eb8\u6807\u6ce8\u793a\u4f8b\u3011')

# ── SLIDE 7 Single check details ──────────────────────────────
s = prs.slides.add_slide(BLANK)
cbg(s)
hdr(s, u'\u5355\u5f20\u68c0\u67e5\u2014\u2014\u8be6\u7ec6\u6b65\u9aa4')
s7 = [
    u'\u2460 \u5728\u63d2\u4ef6\u7a97\u53e3\u4e2d\uff0c\u70b9\u51fb\u5de6\u4fa7\u300c\u5355\u5f20\u68c0\u67e5\u300d\u9009\u9879\u5361\u3002',
    u'\u2461 \u70b9\u51fb\u300c\u68c0\u67e5\u5f53\u524d\u56fe\u7eb8\u300d\u6309\u9215\uff0c\u7cfb\u7edf\u5f00\u59cb\u81ea\u52a8\u68c0\u67e5\u5f53\u524d\u5df2\u6253\u5f00\u7684\u56fe\u7eb8\u3002',
    u'\u2462 \u68c0\u67e5\u5b8c\u6210\u540e\uff1a',
    u'     \u2022 \u5f39\u7a97\u5c06\u5c55\u793a\u68c0\u67e5\u6458\u8981\uff08\u4ec5\u4f9b\u53c2\u8003\uff09\u3002',
    u'     \u2022 \u56fe\u7eb8\u5185\u5c06\u5728\u9519\u8bef\u4f4d\u7f6e\u81ea\u52a8\u6dfb\u52a0\u6807\u6ce8\uff08\u4ee5\u56fe\u7eb8\u6807\u6ce8\u4e3a\u51c6\uff09\u3002',
    u'\u2463 \u4ed4\u7ec6\u6838\u5bf9\u56fe\u7eb8\u5185\u7684\u6807\u6ce8\u5185\u5bb9\uff0c\u786e\u8ba4\u9700\u8981\u4fee\u6539\u7684\u9879\u76ee\u3002',
]
steps(s, s7, top=Inches(0.9), sz=Pt(15))
callout(s, u'[NOTE] \u91cd\u8981\u8bf4\u660e\uff1a\u5f39\u7a97\u4e2d\u7684\u68c0\u67e5\u7ed3\u679c\u4e3a\u8f85\u52a9\u53c2\u8003\uff0c\u56fe\u7eb8\u5185\u7684\u6807\u6ce8\u662f\u6700\u7ec8\u51c6\u786e\u4f9d\u636e\uff0c\u8bf7\u4ee5\u56fe\u7eb8\u6807\u6ce8\u4e3a\u51c6\u3002',
        Inches(0.6), Inches(4.25), Inches(11.8), Inches(0.75),
        bg=BLUE_MID, fg=WHITE, sz=Pt(13))
ph(s, Inches(0.6), Inches(5.2), Inches(11.8), Inches(1.9),
   u'\u3010\u622a\u56fe\u5360\u4f4d\uff1a\u68c0\u67e5\u5f39\u7a97 & \u56fe\u7eb8\u5185\u6807\u6ce8\u5bf9\u6bd4\u793a\u610f\u3011')

# ── SLIDE 8 Edit and recheck ──────────────────────────────────
s = prs.slides.add_slide(BLANK)
cbg(s)
hdr(s, u'\u5355\u5f20\u68c0\u67e5\u2014\u2014\u7f16\u8f91\u4e0e\u590d\u67e5')
s8 = [
    u'\u2460 \u5982\u9700\u4fee\u6539\u56fe\u7eb8\u5185\u5bb9\uff0c\u5148\u5728\u63d2\u4ef6\u7a97\u53e3\u4e2d\u70b9\u51fb\u300c\u5173\u95ed\u300d\u6309\u9215\u3002',
    u'\u2461 \u5728 AutoCAD \u4e2d\u5b8c\u6210\u56fe\u7eb8\u5185\u5bb9\u7684\u4fee\u6539\u3002',
    u'\u2462 \u4fee\u6539\u5b8c\u6210\u540e\uff0c\u91cd\u65b0\u8fd0\u884c CHECKDRAWING\uff0c\u518d\u6b21\u70b9\u51fb\u300c\u68c0\u67e5\u5f53\u524d\u56fe\u7eb8\u300d\uff0c\u8fdb\u884c\u590d\u67e5\u3002',
    u'\u2463 \u82e5\u8981\u6e05\u9664\u4e0a\u4e00\u6b21\u68c0\u67e5\u7559\u4e0b\u7684\u6807\u6ce8\uff0c\u70b9\u51fb\u300c\u6e05\u9664\u5f53\u524d\u56fe\u7eb8\u4fee\u6539\u6ce8\u91ca\u300d\u3002',
    u'   \u2714  \u6e05\u9664\u540e\uff0c\u56fe\u7eb8\u5185\u7684\u6240\u6709\u68c0\u67e5\u6807\u6ce8\u5c06\u88ab\u5220\u9664\uff0c\u53ef\u91cd\u65b0\u5f00\u59cb\u68c0\u67e5\u3002',
]
steps(s, s8, top=Inches(0.95), sz=Pt(15))
flow2 = [u'\u5173\u95ed\u63d2\u4ef6', u'\u4fee\u6539\u56fe\u7eb8', u'\u91cd\u65b0\u68c0\u67e5', u'\uff08\u53ef\u9009\uff09\u6e05\u9664\u6ce8\u91ca']
for i, label in enumerate(flow2):
    lx = Inches(0.6 + i*3.18)
    add_rect(s, lx, Inches(5.2), Inches(2.6), Inches(0.8), fill=BLUE_DARK)
    tb(s, label, lx, Inches(5.2), Inches(2.6), Inches(0.8),
       sz=Pt(14), bold=True, clr=WHITE, align=PP_ALIGN.CENTER)
    if i < 3:
        tb(s, u'\u25b6', lx+Inches(2.6), Inches(5.4), Inches(0.55), Inches(0.4),
           sz=Pt(16), clr=BLUE_MID, align=PP_ALIGN.CENTER)
ph(s, Inches(0.6), Inches(6.1), Inches(11.8), Inches(1.1),
   u'\u3010\u622a\u56fe\u5360\u4f4d\uff1a\u300c\u6e05\u9664\u5f53\u524d\u56fe\u7eb8\u4fee\u6539\u6ce8\u91ca\u300d\u6309\u9215\u4f4d\u7f6e\u3011')

# ── SLIDE 9 Batch check ───────────────────────────────────────
s = prs.slides.add_slide(BLANK)
cbg(s)
hdr(s, u'\u6279\u91cf\u68c0\u67e5\u2014\u2014\u6267\u884c\u6d41\u7a0b')
s9 = [
    u'\u2460 \u5728\u63d2\u4ef6\u7a97\u53e3\u4e2d\uff0c\u70b9\u51fb\u300c\u6279\u91cf\u68c0\u67e5\u300d\u9009\u9879\u5361\u3002',
    u'\u2461 \u70b9\u51fb\u300c\u6267\u884c\u6279\u91cf\u68c0\u67e5\u300d\u6309\u9215\u3002',
    u'\u2462 \u5728\u5f39\u51fa\u7684\u6587\u4ef6\u5939\u9009\u62e9\u5bf9\u8bdd\u6846\u4e2d\uff0c\u9009\u62e9\u9700\u8981\u68c0\u67e5\u7684\u76ee\u6807\u6587\u4ef6\u5939\uff08\u53ef\u9009\u8ba1\u7b97\u673a\u4e0a\u4efb\u610f\u4f4d\u7f6e\uff09\u3002',
    u'\u2463 \u786e\u8ba4\u9009\u62e9\u540e\uff0c\u7cfb\u7edf\u5c06\u81ea\u52a8\u904d\u5386\u5e76\u68c0\u67e5\u8be5\u6587\u4ef6\u5939\u5185\u7684\u6240\u6709 .dwg \u56fe\u7eb8\u6587\u4ef6\u3002',
    u'\u2464 \u7b49\u5f85\u68c0\u67e5\u5b8c\u6210\uff08\u6587\u4ef6\u8f83\u591a\u65f6\u53ef\u80fd\u9700\u8981\u51e0\u5206\u949f\uff0c\u8bf7\u8010\u5fc3\u7b49\u5f85\uff09\u3002',
]
steps(s, s9, top=Inches(0.95), sz=Pt(15))
callout(s, u'\u26a0  \u6ce8\u610f\uff1a\u6267\u884c\u6279\u91cf\u68c0\u67e5\u671f\u95f4\uff0c\u6240\u6709\u76ee\u6807\u56fe\u7eb8\u6587\u4ef6\u5fc5\u987b\u5904\u4e8e\u5173\u95ed\u72b6\u6001\uff0c\u5426\u5219\u53ef\u80fd\u5bfc\u81f4\u68c0\u67e5\u5931\u8d25\u3002',
        Inches(0.6), Inches(4.6), Inches(11.8), Inches(0.7), sz=Pt(13))
ph(s, Inches(0.6), Inches(5.45), Inches(5.6), Inches(1.75),
   u'\u3010\u622a\u56fe\u5360\u4f4d\uff1a\u6279\u91cf\u68c0\u67e5\u6309\u9215 & \u6267\u884c\u6279\u91cf\u68c0\u67e5\u3011')
ph(s, Inches(6.5), Inches(5.45), Inches(6.2), Inches(1.75),
   u'\u3010\u622a\u56fe\u5360\u4f4d\uff1a\u6587\u4ef6\u5939\u9009\u62e9\u5bf9\u8bdd\u6846\u3011')

# ── SLIDE 10 Batch report ─────────────────────────────────────
s = prs.slides.add_slide(BLANK)
cbg(s)
hdr(s, u'\u6279\u91cf\u68c0\u67e5\u2014\u2014\u67e5\u770b\u62a5\u544a')
s10 = [
    u'\u5b8c\u6210\u540e\u4f1a\u5f39\u51fa\u63d0\u793a\u6846\uff1a',
    u'   \u2022 \u70b9\u51fb [\u662f] \u2192 \u81ea\u52a8\u6253\u5f00\u672c\u6b21\u68c0\u67e5\u62a5\u544a\u3002',
    u'   \u2022 \u70b9\u51fb [\u5426] \u2192 \u7a0d\u540e\u53ef\u901a\u8fc7\u4ee5\u4e0b\u65b9\u5f0f\u624b\u52a8\u6253\u5f00\uff1a',
    u'       \u2460 \u5728\u63d2\u4ef6\u7a97\u53e3\u4e2d\u70b9\u51fb\u300c\u6253\u5f00\u6279\u91cf\u68c0\u67e5\u62a5\u544a\u300d\u6309\u9215\uff1b',
    u'       \u2461 \u6216\u5728\u68c0\u67e5\u7684\u76ee\u6807\u6587\u4ef6\u5939\u4e2d\u76f4\u63a5\u627e\u5230\u62a5\u544a\u6587\u4ef6\uff0c\u53cc\u51fb\u6253\u5f00\u3002',
    u'',
    u'\u6253\u5f00\u62a5\u544a\u540e\uff1a',
    u'   \u2022 \u62a5\u544a\u4ee5\u8868\u683c\u5f62\u5f0f\u5217\u51fa\u6bcf\u5f20\u56fe\u7eb8\u53ca\u5176\u5bf9\u5e94\u9519\u8bef\u9879\u3002',
    u'   \u2022 \u70b9\u51fb\u62a5\u544a\u4e2d\u7684\u300c\u6253\u5f00\u56fe\u7eb8\u300d\u6309\u9215\uff0c\u53ef\u76f4\u63a5\u5728 AutoCAD \u4e2d\u6253\u5f00\u5bf9\u5e94 .dwg \u6587\u4ef6\uff0c\u67e5\u770b\u56fe\u7eb8\u5185\u7684\u81ea\u52a8\u6807\u6ce8\u3002',
]
steps(s, s10, top=Inches(0.95), sz=Pt(13.5))
ph(s, Inches(0.6), Inches(5.85), Inches(5.6), Inches(1.4),
   u'\u3010\u622a\u56fe\u5360\u4f4d\uff1a\u5b8c\u6210\u63d0\u793a\u5f39\u7a97\uff08\u662f/\u5426\uff09\u3011')
ph(s, Inches(6.5), Inches(5.85), Inches(6.2), Inches(1.4),
   u'\u3010\u622a\u56fe\u5360\u4f4d\uff1a\u6279\u91cf\u68c0\u67e5\u62a5\u544a\u8868\u683c\u793a\u4f8b\u3011')

# ── SLIDE 11 Batch clear ──────────────────────────────────────
s = prs.slides.add_slide(BLANK)
cbg(s)
hdr(s, u'\u6279\u91cf\u6ce8\u91ca\u6e05\u9664')
s11 = [
    u'\u9002\u7528\u573a\u666f\uff1a\u9700\u8981\u4e00\u6b21\u6027\u6e05\u9664\u67d0\u4e2a\u6587\u4ef6\u5939\u5185\u6240\u6709\u56fe\u7eb8\u7684\u68c0\u67e5\u6807\u6ce8\u65f6\u4f7f\u7528\u3002',
    u'',
    u'\u2460 \u5728\u63d2\u4ef6\u7a97\u53e3\u4e2d\uff0c\u70b9\u51fb\u300c\u6e05\u9664\u6240\u6709\u56fe\u7eb8\u4fee\u6539\u6ce8\u91ca\u300d\u6309\u9215\u3002',
    u'\u2461 \u5728\u5f39\u51fa\u7684\u6587\u4ef6\u5939\u9009\u62e9\u5bf9\u8bdd\u6846\u4e2d\uff0c\u9009\u62e9\u76ee\u6807\u6587\u4ef6\u5939\u3002',
    u'\u2462 \u7cfb\u7edf\u5c06\u81ea\u52a8\u6e05\u9664\u8be5\u6587\u4ef6\u5939\u5185\u6240\u6709 .dwg \u56fe\u7eb8\u4e2d\u7684\u68c0\u67e5\u6807\u6ce8\u3002',
    u'\u2463 \u5b8c\u6210\u540e\uff0c\u76ee\u6807\u56fe\u7eb8\u5c06\u6062\u590d\u4e3a\u65e0\u6807\u6ce8\u72b6\u6001\u3002',
]
steps(s, s11, top=Inches(0.95), sz=Pt(15))
callout(s, u'\u26a0  \u6ce8\u610f\uff1a\u6e05\u9664\u64cd\u4f5c\u524d\uff0c\u6240\u6709\u76ee\u6807\u56fe\u7eb8\u6587\u4ef6\u5fc5\u987b\u5904\u4e8e\u5173\u95ed\u72b6\u6001\uff0c\u5426\u5219\u53ef\u80fd\u5bfc\u81f4\u64cd\u4f5c\u5931\u8d25\u3002',
        Inches(0.6), Inches(4.6), Inches(11.8), Inches(0.7), sz=Pt(13))
ph(s, Inches(0.6), Inches(5.45), Inches(11.8), Inches(1.75),
   u'\u3010\u622a\u56fe\u5360\u4f4d\uff1a\u300c\u6e05\u9664\u6240\u6709\u56fe\u7eb8\u4fee\u6539\u6ce8\u91ca\u300d\u6309\u9215 & \u6587\u4ef6\u5939\u9009\u62e9\u3011')

# ── SLIDE 12 Inspection items ─────────────────────────────────
s = prs.slides.add_slide(BLANK)
cbg(s)
hdr(s, u'\u53ef\u68c0\u67e5\u95ee\u9898\u6c47\u603b')
items = [
    (u'\u9879\u76ee\u53f7\u4e0e\u6587\u4ef6\u540d\u4e0d\u5339\u914d',
     u'\u68c0\u67e5\u56fe\u7eb8\u5185\u662f\u5426\u5b58\u5728\u4e0e\u6587\u4ef6\u540d\u4e0d\u7b26\u7684\u9879\u76ee\u53f7\n\u2192 \u5728\u56fe\u7eb8\u5185\u9519\u8bef\u4f4d\u7f6e\u65c1\u6807\u6ce8\u53c2\u8003\u6b63\u786e\u5185\u5bb9'),
    (u'\u6807\u9898\u680f/BOM\u8868\u56fe\u53f7\u4e0e\u6587\u4ef6\u540d\u4e0d\u5339\u914d',
     u'\u9a8c\u8bc1\u6807\u9898\u680f\u6216 BOM \u8868\u4e2d\u7684\u56fe\u53f7\u662f\u5426\u4e0e\u6587\u4ef6\u540d\u4e00\u81f4\n\u2192 \u5728\u56fe\u7eb8\u5185\u9519\u8bef\u56fe\u53f7\u65c1\u6807\u6ce8\u53c2\u8003\u6b63\u786e\u5185\u5bb9'),
    (u'BOM\u8868\u96f6\u4ef6\u4e0e\u6807\u51c6\u5e93\u6bd4\u5bf9',
     u'\u5bf9\u7167\u6807\u51c6\u96f6\u4ef6\u5e93\uff0c\u68c0\u67e5 BOM \u8868\u4e2d\u96f6\u4ef6\u7684\u5185\u5bb9\u6216\u683c\u5f0f\n\u2192 \u8f93\u51fa\u6b63\u786e\u56fe\u53f7/\u540d\u79f0\uff0c\u6216\u63d0\u793a\u201c\u672a\u6536\u5f55\u201d'),
    (u'\u66f4\u6539\u8bb0\u5f55\u7a7a\u7f3a\u68c0\u67e5',
     u'\u68c0\u67e5\u66f4\u6539\u8bb0\u5f55\u680f\u662f\u5426\u5b58\u5728\u5fc5\u586b\u9879\u672a\u586b\u5199\u7684\u60c5\u51b5'),
    (u'\u6807\u9898\u680f\u7a7a\u7f3a\u68c0\u67e5',
     u'\u68c0\u67e5\u6807\u9898\u680f\u5404\u5b57\u6bb5\u662f\u5426\u5b58\u5728\u7a7a\u7f3a\n\u203b \u6750\u6599/\u89c4\u683c/\u8868\u9762\u5904\u7406\u82e5\u672a\u586b\u5199\uff0c\u4e0d\u4f1a\u5728\u56fe\u4e2d\u5708\u51fa\uff0c\u4f46\u4f1a\u5728\u68c0\u67e5\u62a5\u544a\u4e2d\u8bb0\u5f55'),
]
ch = Inches(1.08)
for i, (title, desc) in enumerate(items):
    col = i % 2; row = i // 2
    lx = Inches(0.5) + col*Inches(6.5)
    ty = Inches(0.88) + row*(ch+Inches(0.12))
    add_rect(s, lx, ty, Inches(6.2), ch, fill=BLUE_LIGHT, line=BLUE_MID, lw=Pt(1))
    add_rect(s, lx, ty, Inches(0.12), ch, fill=BLUE_MID)
    tb(s, title, lx+Inches(0.2), ty+Inches(0.04), Inches(5.9), Inches(0.38),
       sz=Pt(13), bold=True, clr=BLUE_DARK)
    tb(s, desc, lx+Inches(0.2), ty+Inches(0.42), Inches(5.9), Inches(0.62),
       sz=Pt(11), clr=TEXT_DARK)

# ── SLIDE 13 Error handling ───────────────────────────────────
s = prs.slides.add_slide(BLANK)
cbg(s)
hdr(s, u'\u62a5\u9519\u4e0e\u9632\u5d29\u6e83\u673a\u5236')
tb(s, u'\u5f53\u68c0\u67e5\u8fc7\u7a0b\u4e2d\u9047\u5230\u4e25\u91cd\u9519\u8bef\u65f6\uff0c\u7cfb\u7edf\u4f1a\u6709\u4ee5\u4e0b\u8868\u73b0\uff1a',
   Inches(0.7), Inches(0.95), Inches(11), Inches(0.45),
   sz=Pt(15), bold=True, clr=BLUE_DARK)
events = [
    (u'1', u'\u68c0\u67e5\u8fdb\u7a0b\u4e2d\u6b62', u'\u5f53\u524d\u68c0\u67e5\u64cd\u4f5c\u5c06\u505c\u6b62\u6267\u884c'),
    (u'2', u'AutoCAD \u62a5\u9519\u5e76\u5173\u95ed', u'AutoCAD \u5e94\u7528\u7a0b\u5e8f\u53ef\u80fd\u81ea\u52a8\u9000\u51fa'),
    (u'3', u'\u539f\u59cb\u6587\u4ef6\u5b89\u5168', u'\u539f\u59cb .dwg \u6587\u4ef6\u4e0d\u4f1a\u88ab\u635f\u574f\u6216\u4fee\u6539'),
]
for i, (num, ttl, desc) in enumerate(events):
    ty = Inches(1.6) + i*Inches(1.15)
    add_rect(s, Inches(0.6), ty, Inches(0.7), Inches(0.8), fill=BLUE_MID)
    tb(s, num, Inches(0.6), ty, Inches(0.7), Inches(0.8),
       sz=Pt(22), bold=True, clr=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(1.4), ty, Inches(11.0), Inches(0.8),
             fill=GRAY_LIGHT, line=BLUE_LIGHT, lw=Pt(1))
    tb(s, u'{}  \u2014  {}'.format(ttl, desc),
       Inches(1.55), ty+Inches(0.15), Inches(10.7), Inches(0.55),
       sz=Pt(14), clr=TEXT_DARK)
callout(s, u'\u2705  \u89e3\u51b3\u65b9\u6cd5\uff1a\u91cd\u542f AutoCAD\uff0c\u7b49\u5f85\u7247\u523b\u540e\u91cd\u65b0\u6267\u884c\u68c0\u67e5\u64cd\u4f5c\u5373\u53ef\u3002\u65e0\u9700\u62c5\u5fc3\u6587\u4ef6\u5b89\u5168\u3002',
        Inches(0.6), Inches(5.3), Inches(11.8), Inches(0.7),
        bg=GREEN, fg=WHITE, sz=Pt(14))
tb(s, u'\u8bf7\u653e\u5fc3\uff1a\u539f\u59cb\u56fe\u7eb8\u6587\u4ef6\u5747\u4e0d\u4f1a\u56e0\u62a5\u9519\u800c\u88ab\u635f\u574f\u3002',
   Inches(0.6), Inches(6.2), Inches(11.8), Inches(0.45),
   sz=Pt(13), clr=BLUE_MID, align=PP_ALIGN.CENTER)

# ── SLIDE 14 Uninstall ────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
cbg(s)
hdr(s, u'\u8f6f\u4ef6\u5378\u8f7d')
s14 = [
    u'\u2460 \u6253\u5f00 Windows\u3010\u63a7\u5236\u9762\u677f\u3011',
    u'\u2461 \u70b9\u51fb\u3010\u7a0b\u5e8f\u3011',
    u'\u2462 \u70b9\u51fb\u3010\u7a0b\u5e8f\u548c\u529f\u80fd\u3011',
    u'\u2463 \u5728\u7a0b\u5e8f\u5217\u8868\u4e2d\u627e\u5230\u3010CADCheckTool\u3011',
    u'\u2464 \u53f3\u952e\u70b9\u51fb \u2192 \u9009\u62e9\u3010\u5378\u8f7d\u3011',
    u'\u2465 \u6309\u7167\u5378\u8f7d\u5411\u5bfc\u5b8c\u6210\u64cd\u4f5c',
]
for i, step in enumerate(s14):
    ty = Inches(0.95) + i*Inches(0.82)
    add_rect(s, Inches(0.6), ty, Inches(8.5), Inches(0.7),
             fill=BLUE_LIGHT, line=BLUE_MID, lw=Pt(1))
    tb(s, step, Inches(0.75), ty+Inches(0.1), Inches(8.2), Inches(0.52),
       sz=Pt(15), clr=TEXT_DARK)
ph(s, Inches(9.3), Inches(0.95), Inches(3.7), Inches(4.9),
   u'\u3010\u622a\u56fe\u5360\u4f4d\uff1a\u63a7\u5236\u9762\u677f \u2192 \u7a0b\u5e8f\u548c\u529f\u80fd \u5378\u8f7d\u754c\u9762\u3011')

# ── SLIDE 15 Checklist & Support ──────────────────────────────
s = prs.slides.add_slide(BLANK)
blue_bg(s)
tb(s, u'\u5feb\u901f\u6838\u67e5\u6e05\u5355  &  \u83b7\u53d6\u652f\u6301',
   Inches(0.5), Inches(0.3), Inches(12), Inches(0.65),
   sz=Pt(22), bold=True, clr=WHITE, align=PP_ALIGN.CENTER)
checklist = [
    u'\u25a1  \u5b89\u88c5\u524d\u5df2\u5173\u95ed AutoCAD',
    u'\u25a1  \u4ee5\u7ba1\u7406\u5458\u8eab\u4efd\u8fd0\u884c setup.exe \u5e76\u70b9\u51fb Install',
    u'\u25a1  \u4f7f\u7528\u547d\u4ee4 CHECKDRAWING \u6253\u5f00\u63d2\u4ef6\u7a97\u53e3',
    u'\u25a1  \u5355\u5f20\u68c0\u67e5\uff1a\u70b9\u51fb [\u5355\u5f20\u68c0\u67e5] \u2192 [\u68c0\u67e5\u5f53\u524d\u56fe\u7eb8]',
    u'\u25a1  \u7f16\u8f91\u524d\u5148\u70b9\u51fb [\u5173\u95ed]\uff1b\u4fee\u6539\u540e\u91cd\u65b0\u68c0\u67e5',
    u'\u25a1  \u6279\u91cf\u68c0\u67e5\u524d\u786e\u8ba4\u6240\u6709\u76ee\u6807\u56fe\u7eb8\u5df2\u5173\u95ed',
    u'\u25a1  \u6279\u91cf\u68c0\u67e5\uff1a[\u6279\u91cf\u68c0\u67e5] \u2192 [\u6267\u884c\u6279\u91cf\u68c0\u67e5] \u2192 \u9009\u6587\u4ef6\u5939',
    u'\u25a1  \u67e5\u770b\u62a5\u544a\uff1a\u70b9\u51fb [\u662f] \u6216\u4f7f\u7528 [\u6253\u5f00\u6279\u91cf\u68c0\u67e5\u62a5\u544a]',
    u'\u25a1  \u6e05\u9664\u6ce8\u91ca\uff1a\u76ee\u6807\u56fe\u7eb8\u5173\u95ed\u540e\u518d\u6267\u884c\u6e05\u9664\u64cd\u4f5c',
]
for i, item in enumerate(checklist):
    col = i % 2; row = i // 2
    lx = Inches(0.6) + col*Inches(6.5)
    ty = Inches(1.15) + row*Inches(0.76)
    add_rect(s, lx, ty, Inches(6.0), Inches(0.62), fill=RGBColor(0x12,0x4A,0x85))
    tb(s, item, lx+Inches(0.12), ty+Inches(0.08), Inches(5.8), Inches(0.48),
       sz=Pt(13), clr=WHITE)
add_rect(s, Inches(0.6), Inches(5.6), Inches(11.8), Inches(1.25),
         fill=RGBColor(0x0A,0x2A,0x50))
tb(s, u'\u5982\u9700\u5e2e\u52a9\u6216\u53cd\u9988\u95ee\u9898\uff0c\u8bf7\u8054\u7cfb\u6280\u672f\u652f\u6301\uff1a',
   Inches(0.8), Inches(5.68), Inches(5), Inches(0.45),
   sz=Pt(13), bold=True, clr=RGBColor(0xA8,0xD4,0xFF))
tb(s, u'[Email]  \u3010\u8054\u7cfb\u90ae\u7b261\u5360\u4f4d\u3011\n[Tel]    \u3010\u8054\u7cfb\u7535\u8bdd\u5360\u4f4d\u3011\n[Web]    \u3010\u5b98\u65b9\u7f51\u5740\u5360\u4f4d\u3011',
   Inches(0.8), Inches(6.12), Inches(5.5), Inches(0.65),
   sz=Pt(12), clr=RGBColor(0xC8,0xE4,0xFF))
tb(s, u'\u611f\u8c22\u4f7f\u7528 CADCheckTool\uff01',
   Inches(6.5), Inches(5.7), Inches(5.5), Inches(0.9),
   sz=Pt(22), bold=True, clr=RGBColor(0xA8,0xD4,0xFF), align=PP_ALIGN.CENTER)

# ── Save ──────────────────────────────────────────────────────
out = '/home/runner/work/CADCheckTool_1/CADCheckTool_1/CADCheckTool_\u7528\u6237\u7aef\u64cd\u4f5c\u624b\u518c.pptx'
prs.save(out)
print('Saved:', out)
print('Size:', os.path.getsize(out), 'bytes', '({:.1f} KB)'.format(os.path.getsize(out)/1024))
